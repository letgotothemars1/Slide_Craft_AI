from __future__ import annotations

import io
import logging
import re
from dataclasses import dataclass
from pathlib import Path
from urllib import error, parse, request

from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas

from app.config import settings
from app.services.llm_service import LayoutType, PresentationSpec, SlideSpec, ThemeVariant, VisualDensity

logger = logging.getLogger(__name__)

SAFE_NAME_PATTERN = re.compile(r"[^A-Za-z0-9_-]")
PDF_SAFE_REGULAR_FONT = "SlideCraftSafeRegular"
PDF_SAFE_BOLD_FONT = "SlideCraftSafeBold"
PDF_FONT_CANDIDATE_PAIRS: list[tuple[Path, Path | None]] = [
    (
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
    ),
    (
        Path("/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf"),
        Path("/usr/share/fonts/truetype/noto/NotoSans-Bold.ttf"),
    ),
    (
        Path("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"),
        Path("/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"),
    ),
    (
        Path("/usr/share/fonts/truetype/freefont/FreeSans.ttf"),
        Path("/usr/share/fonts/truetype/freefont/FreeSansBold.ttf"),
    ),
    (
        Path("/Library/Fonts/Arial Unicode.ttf"),
        Path("/Library/Fonts/Arial Bold.ttf"),
    ),
    (
        Path("/System/Library/Fonts/Supplemental/Arial.ttf"),
        Path("/System/Library/Fonts/Supplemental/Arial Bold.ttf"),
    ),
    (
        Path("/System/Library/Fonts/Supplemental/Arial Unicode.ttf"),
        Path("/System/Library/Fonts/Supplemental/Arial Bold.ttf"),
    ),
]
CYRILLIC_PATTERN = re.compile(r"[А-Яа-яЁё]")
IMAGE_FETCH_CACHE: dict[str, bytes] = {}

# Fallback mapping when LLM did not provide layout_type.
FALLBACK_TEMPLATE_BY_SLIDE_TYPE: dict[str, LayoutType] = {
    "title": "hero_minimal",
    "agenda": "agenda_clean",
    "content": "content_two_column",
    "quote": "quote_focus",
    "conclusion": "comparison_split",
}

DENSITY_SEQUENCE: tuple[VisualDensity, ...] = ("low", "medium", "high")


@dataclass(frozen=True)
class DeckTheme:
    """Palette + composition settings used by PDF/PPTX renderer."""

    name: ThemeVariant
    page_bg: str
    hero_bg: str
    accent: str
    accent_soft: str
    panel_bg: str
    panel_alt_bg: str
    panel_border: str
    text_primary: str
    text_secondary: str
    text_inverse: str
    rail: str
    connector: str


THEMES: dict[ThemeVariant, DeckTheme] = {
    "dark_tech_pitch": DeckTheme(
        name="dark_tech_pitch",
        page_bg="#0B1020",
        hero_bg="#111A30",
        accent="#22C55E",
        accent_soft="#1A2E3B",
        panel_bg="#162238",
        panel_alt_bg="#1E293B",
        panel_border="#334155",
        text_primary="#F8FAFC",
        text_secondary="#A3B2C8",
        text_inverse="#0B1020",
        rail="#22C55E",
        connector="#4B5563",
    ),
    "clean_editorial": DeckTheme(
        name="clean_editorial",
        page_bg="#FFFCF7",
        hero_bg="#F4F1EB",
        accent="#334155",
        accent_soft="#ECE7DD",
        panel_bg="#FFFFFF",
        panel_alt_bg="#F8FAFC",
        panel_border="#D6D3D1",
        text_primary="#111827",
        text_secondary="#57534E",
        text_inverse="#FFFFFF",
        rail="#6B7280",
        connector="#A8A29E",
    ),
    "infographic_bright": DeckTheme(
        name="infographic_bright",
        page_bg="#F0F9FF",
        hero_bg="#E0F2FE",
        accent="#0EA5E9",
        accent_soft="#D1FAE5",
        panel_bg="#FFFFFF",
        panel_alt_bg="#ECFEFF",
        panel_border="#7DD3FC",
        text_primary="#0F172A",
        text_secondary="#0369A1",
        text_inverse="#FFFFFF",
        rail="#0EA5E9",
        connector="#38BDF8",
    ),
}


def sanitize_job_id_for_paths(job_id: str) -> str:
    sanitized = SAFE_NAME_PATTERN.sub("_", job_id).strip("_")
    return sanitized or "job"


def build_artifact_filename(job_id: str, extension: str) -> str:
    safe_job_id = sanitize_job_id_for_paths(job_id)
    return f"{safe_job_id}.{extension}"


def build_storage_key(job_id: str, extension: str) -> str:
    safe_job_id = sanitize_job_id_for_paths(job_id)
    filename = build_artifact_filename(job_id, extension)
    return f"jobs/{safe_job_id}/{filename}"


def _pdf_latin1_safe(text: str) -> str:
    return text.encode("latin-1", errors="replace").decode("latin-1")


def sanitize_to_ascii_or_latin1(text: str) -> str:
    latin1_safe = _pdf_latin1_safe(text)
    try:
        latin1_safe.encode("latin-1")
        return latin1_safe
    except UnicodeEncodeError:
        return latin1_safe.encode("ascii", errors="replace").decode("ascii")


def _contains_cyrillic(text: str | None) -> bool:
    return bool(text and CYRILLIC_PATTERN.search(text))


def _register_pdf_fonts() -> tuple[str, str, bool]:
    """
    Register Cyrillic-safe fonts for all PDF text.
    Returns (regular_name, bold_name, unicode_supported).
    """
    registered = set(pdfmetrics.getRegisteredFontNames())
    if PDF_SAFE_REGULAR_FONT in registered and PDF_SAFE_BOLD_FONT in registered:
        return PDF_SAFE_REGULAR_FONT, PDF_SAFE_BOLD_FONT, True

    for regular_path, bold_path in PDF_FONT_CANDIDATE_PAIRS:
        if not regular_path.exists():
            continue
        try:
            if PDF_SAFE_REGULAR_FONT not in pdfmetrics.getRegisteredFontNames():
                pdfmetrics.registerFont(TTFont(PDF_SAFE_REGULAR_FONT, str(regular_path)))

            if bold_path and bold_path.exists():
                if PDF_SAFE_BOLD_FONT not in pdfmetrics.getRegisteredFontNames():
                    pdfmetrics.registerFont(TTFont(PDF_SAFE_BOLD_FONT, str(bold_path)))
            else:
                # If bold file is unavailable, reuse regular font to stay Unicode-safe.
                if PDF_SAFE_BOLD_FONT not in pdfmetrics.getRegisteredFontNames():
                    pdfmetrics.registerFont(TTFont(PDF_SAFE_BOLD_FONT, str(regular_path)))

            logger.debug(
                "pdf.font.registered regular=%s bold=%s regular_path=%s bold_path=%s",
                PDF_SAFE_REGULAR_FONT,
                PDF_SAFE_BOLD_FONT,
                regular_path,
                bold_path,
            )
            return PDF_SAFE_REGULAR_FONT, PDF_SAFE_BOLD_FONT, True
        except Exception:
            logger.exception("pdf.font.register_failed regular_path=%s bold_path=%s", regular_path, bold_path)

    logger.warning("pdf.font.fallback regular=Helvetica bold=Helvetica-Bold reason=no unicode ttf font")
    return "Helvetica", "Helvetica-Bold", False


def choose_pdf_font(
    text: str | None,
    style: str,
    *,
    regular_font: str,
    bold_font: str,
) -> str:
    """
    Select PDF font with Cyrillic-safe fallback.
    For safety, all Cyrillic text uses known Unicode fonts.
    """
    prefer_bold = style in {"title", "quote", "accent", "kpi", "badge"}
    if _contains_cyrillic(text):
        return bold_font if prefer_bold else regular_font
    # For consistent look and to avoid accidental glyph issues, we still keep safe family.
    return bold_font if prefer_bold else regular_font


def _pdf_text(value: str | None, *, unicode_font: bool) -> str:
    raw = value or ""
    if unicode_font:
        return raw
    return sanitize_to_ascii_or_latin1(raw)


def _wrap_text_for_pdf(
    text: str,
    *,
    font_name: str,
    font_size: int,
    max_width: float,
) -> list[str]:
    words = text.split()
    if not words:
        return []

    lines: list[str] = []
    current = words[0]

    for word in words[1:]:
        candidate = f"{current} {word}"
        if pdfmetrics.stringWidth(candidate, font_name, font_size) <= max_width:
            current = candidate
        else:
            lines.append(current)
            current = word

    lines.append(current)
    return lines


def _draw_wrapped_lines(
    pdf: canvas.Canvas,
    text: str,
    *,
    x: float,
    y: float,
    max_width: float,
    font_name: str,
    font_size: int,
    line_height: float,
    max_lines: int | None = None,
    prefix_first: str = "",
    prefix_other: str = "",
) -> float:
    pdf.setFont(font_name, font_size)
    lines = _wrap_text_for_pdf(text, font_name=font_name, font_size=font_size, max_width=max_width)
    if max_lines is not None:
        lines = lines[:max_lines]

    for idx, line in enumerate(lines):
        prefix = prefix_first if idx == 0 else prefix_other
        pdf.drawString(x, y, f"{prefix}{line}")
        y -= line_height

    return y


def _joined_deck_text(spec: PresentationSpec) -> str:
    slide_text = " ".join(_joined_slide_text(slide) for slide in spec.slides)
    return " ".join(
        [
            spec.title.lower(),
            (spec.subtitle or "").lower(),
            spec.audience.lower(),
            spec.style.lower(),
            slide_text.lower(),
        ]
    )


def _infer_theme_variant(spec: PresentationSpec) -> ThemeVariant:
    text = _joined_deck_text(spec)
    if any(
        word in text
        for word in [
            "ai",
            "product",
            "startup",
            "architecture",
            "technology",
            "pitch",
            "технолог",
            "архитект",
            "продукт",
            "стартап",
        ]
    ):
        return "dark_tech_pitch"
    if any(
        word in text
        for word in [
            "history",
            "culture",
            "academic",
            "conceptual",
            "reflective",
            "истор",
            "культур",
            "академ",
            "концепт",
        ]
    ):
        return "clean_editorial"
    if any(
        word in text
        for word in [
            "education",
            "explainer",
            "framework",
            "method",
            "process",
            "learning",
            "образован",
            "метод",
            "процесс",
            "обуч",
        ]
    ):
        return "infographic_bright"

    # Default fallback by requested style to keep behavior stable.
    if spec.style == "dark":
        return "dark_tech_pitch"
    if spec.style == "creative":
        return "infographic_bright"
    return "clean_editorial"


def _select_deck_theme(spec: PresentationSpec, *, render_format: str) -> DeckTheme:
    if spec.theme_variant in THEMES:
        selected = spec.theme_variant
        source = "spec"
    else:
        selected = _infer_theme_variant(spec)
        source = "fallback"

    logger.debug(
        "theme_variant.selected format=%s theme_variant=%s source=%s",
        render_format,
        selected,
        source,
    )
    return THEMES[selected]


def _select_visual_density(slide: SlideSpec, *, template: LayoutType, render_format: str) -> VisualDensity:
    if slide.visual_density in DENSITY_SEQUENCE:
        density = slide.visual_density
        source = "spec"
    elif template in {"hero_minimal", "quote_focus"}:
        density = "low"
        source = "fallback"
    elif template in {"kpi_cards", "timeline_process", "infographic_visual"}:
        density = "high"
        source = "fallback"
    else:
        density = "medium"
        source = "fallback"

    logger.debug(
        "visual_density.selected format=%s slide_id=%s density=%s source=%s",
        render_format,
        slide.id,
        density,
        source,
    )
    return density


def _density_limit(density: VisualDensity, *, low: int, medium: int, high: int) -> int:
    if density == "low":
        return low
    if density == "high":
        return high
    return medium


def _density_scale(density: VisualDensity) -> float:
    if density == "low":
        return 1.08
    if density == "high":
        return 0.92
    return 1.0


def _pdf_color(value: str) -> colors.Color:
    return colors.HexColor(value)


def _rgb_from_hex(value: str) -> tuple[int, int, int]:
    cleaned = value.lstrip("#")
    return int(cleaned[0:2], 16), int(cleaned[2:4], 16), int(cleaned[4:6], 16)


def _paint_pdf_page_background(pdf: canvas.Canvas, *, width: float, height: float, theme: DeckTheme) -> None:
    pdf.setFillColor(_pdf_color(theme.page_bg))
    pdf.rect(0, 0, width, height, fill=1, stroke=0)


def _select_template(slide: SlideSpec, *, render_format: str) -> LayoutType:
    if slide.layout_type is not None:
        logger.debug(
            "layout_type.selected format=%s slide_id=%s slide_type=%s layout_type=%s source=spec",
            render_format,
            slide.id,
            slide.type,
            slide.layout_type,
        )
        return slide.layout_type

    # Fallback selection by slide semantics when layout_type is missing.
    if slide.type == "title":
        selected = "hero_minimal"
    elif slide.type == "agenda":
        selected = "agenda_clean"
    elif _has_quote_semantics(slide):
        selected = "quote_focus"
    elif _has_comparison_semantics(slide):
        selected = "comparison_split"
    elif _has_timeline_semantics(slide):
        selected = "timeline_process"
    elif _has_metrics_semantics(slide):
        selected = "kpi_cards"
    elif slide.type == "conclusion":
        selected = "hero_minimal"
    else:
        selected = FALLBACK_TEMPLATE_BY_SLIDE_TYPE.get(slide.type, "content_two_column")

    logger.debug(
        "renderer.template.fallback_used format=%s slide_id=%s slide_type=%s selected=%s",
        render_format,
        slide.id,
        slide.type,
        selected,
    )
    logger.debug(
        "layout_type.selected format=%s slide_id=%s slide_type=%s layout_type=%s source=fallback",
        render_format,
        slide.id,
        slide.type,
        selected,
    )
    return selected


def resolve_layout_type(slide: SlideSpec) -> LayoutType:
    """
    Public helper for non-renderer layers (e.g., orchestrator asset planning).
    No renderer-side logging here; logs are emitted in rendering flow.
    """
    if slide.layout_type is not None:
        return slide.layout_type

    if slide.type == "title":
        return "hero_minimal"
    if slide.type == "agenda":
        return "agenda_clean"
    if _has_quote_semantics(slide):
        return "quote_focus"
    if _has_comparison_semantics(slide):
        return "comparison_split"
    if _has_timeline_semantics(slide):
        return "timeline_process"
    if _has_metrics_semantics(slide):
        return "kpi_cards"
    if slide.type == "conclusion":
        return "hero_minimal"
    return FALLBACK_TEMPLATE_BY_SLIDE_TYPE.get(slide.type, "content_two_column")


def _extract_cards(slide: SlideSpec, *, max_items: int = 4) -> list[tuple[str, str]]:
    raw_items = slide.bullets[:]
    if not raw_items and slide.body:
        raw_items = [item.strip() for item in re.split(r"[\n;]", slide.body) if item.strip()]

    cards: list[tuple[str, str]] = []
    for idx, item in enumerate(raw_items[:max_items], start=1):
        match = re.match(r"^\s*([\d][\d\s\.,%kKmMbBxX+-]*)\s*[-:|]?\s*(.*)$", item)
        if match:
            metric = match.group(1).strip()
            label = match.group(2).strip() or f"Metric {idx}"
        else:
            metric = str(idx)
            label = item
        cards.append((metric, label))

    if not cards:
        cards = [("1", "No KPI provided"), ("2", "Add bullets in spec")]

    return cards


def _extract_timeline_steps(slide: SlideSpec, *, max_items: int = 6) -> list[str]:
    steps = [item.strip() for item in slide.bullets if item.strip()]
    if not steps and slide.body:
        steps = [item.strip() for item in re.split(r"[\n\.]", slide.body) if item.strip()]

    if not steps:
        steps = ["Step 1", "Step 2", "Step 3"]

    return steps[:max_items]


def _extract_infographic_nodes(slide: SlideSpec, *, max_items: int = 6) -> list[str]:
    nodes = [item.strip() for item in slide.bullets if item.strip()]
    if not nodes and slide.body:
        nodes = [item.strip() for item in re.split(r"[\n;,.]", slide.body) if item.strip()]

    if not nodes:
        nodes = ["Node A", "Node B", "Node C"]

    return nodes[:max_items]


def _content_variant(slide: SlideSpec) -> str:
    """Simple deterministic variant to avoid repetitive content slide look."""
    seed = sum(ord(ch) for ch in slide.id)
    return "left_text" if seed % 2 == 0 else "left_visual"


def _extract_quote_text(slide: SlideSpec) -> str:
    if slide.body:
        return slide.body
    if slide.key_message:
        return slide.key_message
    if slide.bullets:
        return slide.bullets[0]
    if slide.subtitle:
        return slide.subtitle
    return "Quote placeholder"


def _extract_comparison_columns(
    slide: SlideSpec,
    *,
    max_items_per_side: int = 4,
) -> tuple[str, list[str], str, list[str]]:
    """Split content into two comparison columns."""
    left_title = "Option A"
    right_title = "Option B"

    if slide.section and "/" in slide.section:
        parts = [part.strip() for part in slide.section.split("/", maxsplit=1)]
        if len(parts) == 2:
            left_title, right_title = parts[0], parts[1]

    left_items: list[str] = []
    right_items: list[str] = []

    for item in slide.bullets:
        if "|" in item:
            left, right = item.split("|", maxsplit=1)
            left_items.append(left.strip())
            right_items.append(right.strip())
        else:
            if len(left_items) <= len(right_items):
                left_items.append(item)
            else:
                right_items.append(item)

    if not left_items and not right_items and slide.body:
        chunks = [part.strip() for part in re.split(r"[\n;]", slide.body) if part.strip()]
        midpoint = max(1, len(chunks) // 2)
        left_items = chunks[:midpoint]
        right_items = chunks[midpoint:]

    if not left_items:
        left_items = ["Point A1", "Point A2"]
    if not right_items:
        right_items = ["Point B1", "Point B2"]

    return left_title, left_items[:max_items_per_side], right_title, right_items[:max_items_per_side]


def _joined_slide_text(slide: SlideSpec) -> str:
    parts = [
        slide.title,
        slide.subtitle,
        slide.body,
        slide.visual_hint,
        slide.key_message,
        slide.section,
        slide.chart_hint,
        slide.image_prompt,
        " ".join(slide.bullets),
    ]
    return " ".join([item for item in parts if item]).lower()


def _has_quote_semantics(slide: SlideSpec) -> bool:
    if slide.type == "quote":
        return True
    quote_text = _extract_quote_text(slide)
    return any(mark in quote_text for mark in ['"', "“", "”", "«", "»"])


def _has_comparison_semantics(slide: SlideSpec) -> bool:
    text = _joined_slide_text(slide)
    if any("|" in item for item in slide.bullets):
        return True
    keywords = [
        " vs ",
        "versus",
        "before",
        "after",
        "compare",
        "comparison",
        "до",
        "после",
        "против",
        "сравнение",
    ]
    return any(word in text for word in keywords)


def _has_timeline_semantics(slide: SlideSpec) -> bool:
    text = _joined_slide_text(slide)
    keywords = [
        "timeline",
        "roadmap",
        "process",
        "milestone",
        "step",
        "q1",
        "q2",
        "q3",
        "q4",
        "этап",
        "шаг",
        "дорожн",
    ]
    return any(word in text for word in keywords) or len(slide.bullets) >= 4


def _has_metrics_semantics(slide: SlideSpec) -> bool:
    text = _joined_slide_text(slide)
    numeric_tokens = len(re.findall(r"\d", text))
    metric_words = ["kpi", "metric", "metrics", "roi", "conversion", "%", "рост", "метрик"]
    return numeric_tokens >= 4 or any(word in text for word in metric_words)


def _read_image_bytes(image_url: str | None) -> bytes | None:
    """Load image bytes from URL with small in-memory cache."""
    if not image_url:
        return None

    if image_url in IMAGE_FETCH_CACHE:
        return IMAGE_FETCH_CACHE[image_url]

    # Fast-path for local storage mode (/files/...) to avoid HTTP round-trip.
    try:
        parsed = parse.urlparse(image_url)
        if parsed.path.startswith("/files/"):
            local_rel = parse.unquote(parsed.path.removeprefix("/files/"))
            local_path = settings.STORAGE_PATH / local_rel
            if local_path.exists():
                data = local_path.read_bytes()
                IMAGE_FETCH_CACHE[image_url] = data
                return data
    except Exception:
        logger.exception("renderer.image.fetch_failed source=local_path image_url=%s", image_url)

    try:
        with request.urlopen(image_url, timeout=20) as resp:
            data = resp.read()
        IMAGE_FETCH_CACHE[image_url] = data
        return data
    except error.HTTPError:
        logger.exception("renderer.image.fetch_failed source=http image_url=%s", image_url)
    except Exception:
        logger.exception("renderer.image.fetch_failed source=network image_url=%s", image_url)

    return None


def _draw_pdf_visual_placeholder(
    pdf: canvas.Canvas,
    *,
    slide: SlideSpec,
    theme: DeckTheme,
    density: VisualDensity,
    x: float,
    y: float,
    width: float,
    height: float,
    font_name: str,
    unicode_font: bool,
) -> None:
    """Draw real image when available, otherwise render visual placeholder block."""
    hint = slide.visual_hint or slide.image_prompt or slide.chart_hint or "Visual placeholder"

    # 1) Try real image.
    image_bytes = _read_image_bytes(slide.image_url)
    if image_bytes:
        try:
            img = ImageReader(io.BytesIO(image_bytes))
            iw, ih = img.getSize()
            ratio = min(width / iw, height / ih)
            draw_w = iw * ratio
            draw_h = ih * ratio
            draw_x = x + (width - draw_w) / 2
            draw_y = y + (height - draw_h) / 2
            pdf.drawImage(img, draw_x, draw_y, draw_w, draw_h, preserveAspectRatio=True, anchor="c")
            logger.debug("renderer.real_image.used format=pdf slide_id=%s", slide.id)
            logger.debug(
                "renderer.visual_block.created format=pdf slide_id=%s kind=real_image",
                slide.id,
            )
            return
        except Exception:
            logger.exception("renderer.image.decode_failed format=pdf slide_id=%s", slide.id)

    # 2) Placeholder fallback.
    pdf.setFillColor(_pdf_color(theme.panel_alt_bg))
    pdf.setStrokeColor(_pdf_color(theme.panel_border))
    pdf.roundRect(x, y, width, height, 10, fill=1, stroke=1)

    badge_label = "Image placeholder" if slide.image_prompt else "Visual block"
    logger.debug("renderer.placeholder_image.used format=pdf slide_id=%s", slide.id)
    logger.debug(
        "renderer.visual_block.created format=pdf slide_id=%s kind=%s",
        slide.id,
        badge_label,
    )

    label_font = choose_pdf_font(
        badge_label,
        "badge",
        regular_font=font_name,
        bold_font=PDF_SAFE_BOLD_FONT if PDF_SAFE_BOLD_FONT in pdfmetrics.getRegisteredFontNames() else font_name,
    )
    hint_font = choose_pdf_font(
        hint,
        "body",
        regular_font=font_name,
        bold_font=PDF_SAFE_BOLD_FONT if PDF_SAFE_BOLD_FONT in pdfmetrics.getRegisteredFontNames() else font_name,
    )

    pdf.setFillColor(_pdf_color(theme.text_secondary))
    pdf.setFont(label_font, 10)
    pdf.drawString(x + 12, y + height - 18, badge_label)

    pdf.setFillColor(_pdf_color(theme.text_secondary))
    _draw_wrapped_lines(
        pdf,
        _pdf_text(hint, unicode_font=unicode_font),
        x=x + 12,
        y=y + height - 36,
        max_width=width - 24,
        font_name=hint_font,
        font_size=11,
        line_height=14,
        max_lines=_density_limit(density, low=3, medium=5, high=6),
    )


def _render_pdf_hero_minimal(
    pdf: canvas.Canvas,
    *,
    width: float,
    height: float,
    slide: SlideSpec,
    theme: DeckTheme,
    density: VisualDensity,
    font_name: str,
    unicode_font: bool,
) -> None:
    scale = _density_scale(density)

    pdf.setFillColor(_pdf_color(theme.hero_bg))
    pdf.rect(0, height * 0.62, width, height * 0.38, fill=1, stroke=0)

    pdf.setFillColor(_pdf_color(theme.accent))
    pdf.rect(52, 56, 8, height - 112, fill=1, stroke=0)

    pdf.setFillColor(_pdf_color(theme.accent_soft))
    pdf.roundRect(width - 288, 54, 232, 132, 16, fill=1, stroke=0)

    pdf.setFillColor(_pdf_color(theme.text_primary))
    y = height - 120
    y = _draw_wrapped_lines(
        pdf,
        _pdf_text(slide.title, unicode_font=unicode_font),
        x=84,
        y=y,
        max_width=width - 150,
        font_name=font_name,
        font_size=int(44 * scale),
        line_height=int(50 * scale),
        max_lines=_density_limit(density, low=2, medium=3, high=4),
    )

    if slide.subtitle:
        pdf.setFillColor(_pdf_color(theme.text_secondary))
        y = _draw_wrapped_lines(
            pdf,
            _pdf_text(slide.subtitle, unicode_font=unicode_font),
            x=84,
            y=y - 12,
            max_width=width - 170,
            font_name=font_name,
            font_size=int(18 * scale),
            line_height=int(24 * scale),
            max_lines=3,
        )

    if slide.key_message:
        pdf.setFillColor(_pdf_color(theme.text_secondary))
        pdf.setFont(font_name, 14)
        _draw_wrapped_lines(
            pdf,
            _pdf_text(slide.key_message, unicode_font=unicode_font),
            x=84,
            y=108,
            max_width=width - 170,
            font_name=font_name,
            font_size=14,
            line_height=18,
            max_lines=2,
        )

    if slide.image_url or slide.image_prompt:
        _draw_pdf_visual_placeholder(
            pdf,
            slide=slide,
            theme=theme,
            density=density,
            x=width - 260,
            y=64,
            width=196,
            height=112,
            font_name=font_name,
            unicode_font=unicode_font,
        )


def _render_pdf_agenda_clean(
    pdf: canvas.Canvas,
    *,
    width: float,
    height: float,
    slide: SlideSpec,
    theme: DeckTheme,
    density: VisualDensity,
    font_name: str,
    unicode_font: bool,
) -> None:
    scale = _density_scale(density)
    pdf.setFillColor(_pdf_color(theme.text_primary))
    pdf.setFont(font_name, 32)
    pdf.drawString(56, height - 84, _pdf_text(slide.title, unicode_font=unicode_font))

    if slide.subtitle:
        pdf.setFillColor(_pdf_color(theme.text_secondary))
        pdf.setFont(font_name, 14)
        _draw_wrapped_lines(
            pdf,
            _pdf_text(slide.subtitle, unicode_font=unicode_font),
            x=56,
            y=height - 112,
            max_width=width - 112,
            font_name=font_name,
            font_size=14,
            line_height=19,
            max_lines=2,
        )

    rail_x = 82
    y_top = height - 160
    y_bottom = 90
    pdf.setStrokeColor(_pdf_color(theme.rail))
    pdf.setLineWidth(2)
    pdf.line(rail_x, y_bottom, rail_x, y_top)

    items = slide.bullets[: _density_limit(density, low=4, medium=6, high=7)]
    if not items and slide.body:
        items = [
            item.strip()
            for item in re.split(r"[\n;]", slide.body)
            if item.strip()
        ][:_density_limit(density, low=4, medium=6, high=7)]
    if not items:
        items = ["Agenda item 1", "Agenda item 2", "Agenda item 3"]

    gap = min(80, (y_top - y_bottom) / max(1, len(items) - 1))
    for idx, item in enumerate(items):
        y = y_top - idx * gap
        pdf.setFillColor(_pdf_color(theme.accent))
        pdf.circle(rail_x, y, 7, stroke=0, fill=1)

        badge_x = 106
        badge_y = y - 10
        pdf.setFillColor(_pdf_color(theme.panel_alt_bg))
        pdf.setStrokeColor(_pdf_color(theme.panel_border))
        pdf.roundRect(badge_x, badge_y, 34, 20, 6, fill=1, stroke=1)
        pdf.setFillColor(_pdf_color(theme.text_primary))
        pdf.setFont(font_name, 10)
        pdf.drawCentredString(badge_x + 17, badge_y + 6, str(idx + 1))

        pdf.setFillColor(_pdf_color(theme.text_primary))
        _draw_wrapped_lines(
            pdf,
            _pdf_text(item, unicode_font=unicode_font),
            x=148,
            y=y + 6,
            max_width=width - 210,
            font_name=font_name,
            font_size=int(14 * scale),
            line_height=int(18 * scale),
            max_lines=_density_limit(density, low=2, medium=2, high=3),
        )


def _render_pdf_content_two_column(
    pdf: canvas.Canvas,
    *,
    width: float,
    height: float,
    slide: SlideSpec,
    theme: DeckTheme,
    density: VisualDensity,
    font_name: str,
    unicode_font: bool,
) -> None:
    variant = _content_variant(slide)
    scale = _density_scale(density)

    pdf.setFillColor(_pdf_color(theme.text_primary))
    pdf.setFont(font_name, int(30 * scale))
    pdf.drawString(56, height - 82, _pdf_text(slide.title, unicode_font=unicode_font))

    if slide.subtitle:
        pdf.setFillColor(_pdf_color(theme.text_secondary))
        pdf.setFont(font_name, int(13 * scale))
        pdf.drawString(56, height - 104, _pdf_text(slide.subtitle, unicode_font=unicode_font))

    left_x = 56
    right_x = width * 0.52
    left_w = width * 0.42
    right_w = width - right_x - 56

    if variant == "left_text":
        text_x, text_w = left_x, left_w
        visual_x, visual_w = right_x, right_w
    else:
        text_x, text_w = right_x, right_w
        visual_x, visual_w = left_x, left_w

    # Text block.
    pdf.setFillColor(_pdf_color(theme.text_primary))
    text_y = height - 150
    items = slide.bullets[:]
    if not items and slide.body:
        items = [slide.body]

    if items:
        for item in items[: _density_limit(density, low=4, medium=6, high=8)]:
            text_y = _draw_wrapped_lines(
                pdf,
                _pdf_text(item, unicode_font=unicode_font),
                x=text_x,
                y=text_y,
                max_width=text_w,
                font_name=font_name,
                font_size=int(13 * scale),
                line_height=int(17 * scale),
                max_lines=_density_limit(density, low=2, medium=3, high=4),
                prefix_first="- " if item in slide.bullets else "",
            )
            text_y -= 6

    # Visual panel — only when image is available or explicitly requested.
    if slide.image_url or slide.image_prompt:
        _draw_pdf_visual_placeholder(
            pdf,
            slide=slide,
            theme=theme,
            density=density,
            x=visual_x,
            y=120,
            width=visual_w,
            height=height - 270,
            font_name=font_name,
            unicode_font=unicode_font,
        )

    if slide.key_message:
        pdf.setFillColor(_pdf_color(theme.text_secondary))
        pdf.setFont(font_name, 12)
        _draw_wrapped_lines(
            pdf,
            _pdf_text(slide.key_message, unicode_font=unicode_font),
            x=56,
            y=86,
            max_width=width - 112,
            font_name=font_name,
            font_size=12,
            line_height=15,
            max_lines=2,
        )


def _render_pdf_kpi_cards(
    pdf: canvas.Canvas,
    *,
    width: float,
    height: float,
    slide: SlideSpec,
    theme: DeckTheme,
    density: VisualDensity,
    font_name: str,
    unicode_font: bool,
) -> None:
    scale = _density_scale(density)
    pdf.setFillColor(_pdf_color(theme.text_primary))
    pdf.setFont(font_name, 30)
    pdf.drawString(56, height - 82, _pdf_text(slide.title, unicode_font=unicode_font))

    if slide.key_message:
        pdf.setFillColor(_pdf_color(theme.text_secondary))
        pdf.setFont(font_name, 13)
        _draw_wrapped_lines(
            pdf,
            _pdf_text(slide.key_message, unicode_font=unicode_font),
            x=56,
            y=height - 106,
            max_width=width - 112,
            font_name=font_name,
            font_size=13,
            line_height=17,
            max_lines=2,
        )

    cards = _extract_cards(slide, max_items=_density_limit(density, low=2, medium=4, high=4))
    card_w = (width - 56 - 56 - 20) / 2
    card_h = 138
    start_y = height - 285

    for idx, (metric, label) in enumerate(cards[:4]):
        row = idx // 2
        col = idx % 2
        x = 56 + col * (card_w + 20)
        y = start_y - row * (card_h + 18)

        pdf.setFillColor(_pdf_color(theme.panel_bg if idx % 2 == 0 else theme.panel_alt_bg))
        pdf.setStrokeColor(_pdf_color(theme.panel_border))
        pdf.roundRect(x, y, card_w, card_h, 12, fill=1, stroke=1)

        # KPI badge.
        pdf.setFillColor(_pdf_color(theme.accent))
        pdf.circle(x + card_w - 22, y + card_h - 18, 9, stroke=0, fill=1)
        pdf.setFillColor(_pdf_color(theme.text_inverse))
        pdf.setFont(font_name, 9)
        pdf.drawCentredString(x + card_w - 22, y + card_h - 21, str(idx + 1))

        pdf.setFillColor(_pdf_color(theme.text_primary))
        pdf.setFont(font_name, int(28 * scale))
        _draw_wrapped_lines(
            pdf,
            _pdf_text(metric, unicode_font=unicode_font),
            x=x + 16,
            y=y + 92,
            max_width=card_w - 34,
            font_name=font_name,
            font_size=int(28 * scale),
            line_height=int(30 * scale),
            max_lines=1,
        )

        pdf.setFillColor(_pdf_color(theme.text_secondary))
        pdf.setFont(font_name, 11)
        _draw_wrapped_lines(
            pdf,
            _pdf_text(label, unicode_font=unicode_font),
            x=x + 16,
            y=y + 58,
            max_width=card_w - 30,
            font_name=font_name,
            font_size=11,
            line_height=14,
            max_lines=3,
        )


def _render_pdf_timeline_process(
    pdf: canvas.Canvas,
    *,
    width: float,
    height: float,
    slide: SlideSpec,
    theme: DeckTheme,
    density: VisualDensity,
    font_name: str,
    unicode_font: bool,
) -> None:
    scale = _density_scale(density)
    pdf.setFillColor(_pdf_color(theme.text_primary))
    pdf.setFont(font_name, 30)
    pdf.drawString(56, height - 82, _pdf_text(slide.title, unicode_font=unicode_font))

    if slide.subtitle:
        pdf.setFillColor(_pdf_color(theme.text_secondary))
        pdf.setFont(font_name, int(13 * scale))
        pdf.drawString(56, height - 104, _pdf_text(slide.subtitle, unicode_font=unicode_font))

    steps = _extract_timeline_steps(slide, max_items=_density_limit(density, low=4, medium=6, high=8))
    x_left = 80
    x_right = width - 72
    y = height - 176
    gap = max(58, (height - 280) / max(1, len(steps)))

    pdf.setStrokeColor(_pdf_color(theme.connector))
    pdf.setLineWidth(2)
    pdf.line(x_left, y - (len(steps) - 1) * gap, x_left, y + 6)

    for idx, step in enumerate(steps):
        yy = y - idx * gap
        pdf.setFillColor(_pdf_color(theme.accent))
        pdf.circle(x_left, yy, 8, stroke=0, fill=1)

        pdf.setFillColor(_pdf_color(theme.panel_bg))
        pdf.setStrokeColor(_pdf_color(theme.panel_border))
        pdf.roundRect(100, yy - 20, x_right - 100, 40, 8, fill=1, stroke=1)

        pdf.setFillColor(_pdf_color(theme.text_primary))
        _draw_wrapped_lines(
            pdf,
            _pdf_text(step, unicode_font=unicode_font),
            x=112,
            y=yy + 6,
            max_width=x_right - 140,
            font_name=font_name,
            font_size=int(12 * scale),
            line_height=int(15 * scale),
            max_lines=_density_limit(density, low=2, medium=2, high=3),
            prefix_first=f"{idx + 1}. ",
            prefix_other="   ",
        )

    if slide.visual_hint or slide.image_prompt:
        _draw_pdf_visual_placeholder(
            pdf,
            slide=slide,
            theme=theme,
            density=density,
            x=width - 240,
            y=58,
            width=182,
            height=96,
            font_name=font_name,
            unicode_font=unicode_font,
        )


def _render_pdf_infographic_visual(
    pdf: canvas.Canvas,
    *,
    width: float,
    height: float,
    slide: SlideSpec,
    theme: DeckTheme,
    density: VisualDensity,
    font_name: str,
    unicode_font: bool,
) -> None:
    scale = _density_scale(density)
    pdf.setFillColor(_pdf_color(theme.text_primary))
    pdf.setFont(font_name, 30)
    pdf.drawString(56, height - 82, _pdf_text(slide.title, unicode_font=unicode_font))

    center_x = width / 2
    center_y = height / 2 - 10
    pdf.setFillColor(_pdf_color(theme.accent_soft))
    pdf.setStrokeColor(_pdf_color(theme.panel_border))
    pdf.roundRect(center_x - 118, center_y - 38, 236, 76, 10, fill=1, stroke=1)

    center_text = slide.key_message or slide.body or "Core concept"
    pdf.setFillColor(_pdf_color(theme.text_primary))
    _draw_wrapped_lines(
        pdf,
        _pdf_text(center_text, unicode_font=unicode_font),
        x=center_x - 102,
        y=center_y + 12,
        max_width=204,
        font_name=font_name,
        font_size=int(12 * scale),
        line_height=int(15 * scale),
        max_lines=_density_limit(density, low=2, medium=3, high=4),
    )

    node_texts = _extract_infographic_nodes(slide, max_items=_density_limit(density, low=4, medium=6, high=6))
    node_positions = [
        (center_x - 240, center_y + 108),
        (center_x + 90, center_y + 108),
        (center_x - 240, center_y - 126),
        (center_x + 90, center_y - 126),
        (center_x - 66, center_y + 168),
        (center_x - 66, center_y - 186),
    ]

    for idx, text in enumerate(node_texts):
        x, y = node_positions[idx]
        w, h = 160, 60

        pdf.setStrokeColor(_pdf_color(theme.connector))
        pdf.setLineWidth(1)
        pdf.line(center_x, center_y, x + w / 2, y + h / 2)

        pdf.setFillColor(_pdf_color(theme.panel_bg))
        pdf.setStrokeColor(_pdf_color(theme.panel_border))
        pdf.roundRect(x, y, w, h, 8, fill=1, stroke=1)

        pdf.setFillColor(_pdf_color(theme.text_primary))
        _draw_wrapped_lines(
            pdf,
            _pdf_text(text, unicode_font=unicode_font),
            x=x + 10,
            y=y + 38,
            max_width=w - 20,
            font_name=font_name,
            font_size=10,
            line_height=12,
            max_lines=3,
        )

    _draw_pdf_visual_placeholder(
        pdf,
        slide=slide,
        theme=theme,
        density=density,
        x=56,
        y=52,
        width=186,
        height=76,
        font_name=font_name,
        unicode_font=unicode_font,
    )


def _render_pdf_quote_focus(
    pdf: canvas.Canvas,
    *,
    width: float,
    height: float,
    slide: SlideSpec,
    theme: DeckTheme,
    density: VisualDensity,
    font_name: str,
    unicode_font: bool,
) -> None:
    logger.debug("renderer.quote_layout.used format=pdf slide_id=%s", slide.id)

    quote = _extract_quote_text(slide)
    scale = _density_scale(density)
    pdf.setFillColor(_pdf_color(theme.page_bg))
    pdf.rect(0, 0, width, height, fill=1, stroke=0)

    pdf.setFillColor(_pdf_color(theme.panel_alt_bg))
    pdf.roundRect(42, 44, width - 84, height - 88, 20, fill=1, stroke=0)

    pdf.setFillColor(_pdf_color(theme.accent))
    pdf.setFont(font_name, 72)
    pdf.drawString(72, height - 168, '"')

    pdf.setFillColor(_pdf_color(theme.text_primary))
    _draw_wrapped_lines(
        pdf,
        _pdf_text(quote, unicode_font=unicode_font),
        x=102,
        y=height - 180,
        max_width=width - 180,
        font_name=font_name,
        font_size=int(30 * scale),
        line_height=int(36 * scale),
        max_lines=_density_limit(density, low=7, medium=9, high=11),
    )

    if slide.subtitle:
        pdf.setFillColor(_pdf_color(theme.text_secondary))
        pdf.setFont(font_name, 14)
        pdf.drawString(102, 90, _pdf_text(slide.subtitle, unicode_font=unicode_font))

    if slide.image_url:
        _draw_pdf_visual_placeholder(
            pdf,
            slide=slide,
            theme=theme,
            density=density,
            x=width - 260,
            y=52,
            width=188,
            height=84,
            font_name=font_name,
            unicode_font=unicode_font,
        )


def _render_pdf_comparison_split(
    pdf: canvas.Canvas,
    *,
    width: float,
    height: float,
    slide: SlideSpec,
    theme: DeckTheme,
    density: VisualDensity,
    font_name: str,
    unicode_font: bool,
) -> None:
    logger.debug("renderer.comparison_layout.used format=pdf slide_id=%s", slide.id)
    left_title, left_items, right_title, right_items = _extract_comparison_columns(
        slide,
        max_items_per_side=_density_limit(density, low=3, medium=4, high=5),
    )
    scale = _density_scale(density)

    pdf.setFillColor(_pdf_color(theme.text_primary))
    pdf.setFont(font_name, int(30 * scale))
    pdf.drawString(56, height - 82, _pdf_text(slide.title, unicode_font=unicode_font))

    panel_top = height - 140
    panel_h = height - 220
    gap = 16
    panel_w = (width - 56 - 56 - gap) / 2
    left_x = 56
    right_x = left_x + panel_w + gap

    pdf.setFillColor(_pdf_color(theme.panel_bg))
    pdf.setStrokeColor(_pdf_color(theme.panel_border))
    pdf.roundRect(left_x, 84, panel_w, panel_h, 12, fill=1, stroke=1)

    pdf.setFillColor(_pdf_color(theme.panel_alt_bg))
    pdf.setStrokeColor(_pdf_color(theme.panel_border))
    pdf.roundRect(right_x, 84, panel_w, panel_h, 12, fill=1, stroke=1)

    pdf.setStrokeColor(_pdf_color(theme.connector))
    pdf.setLineWidth(2)
    pdf.line(width / 2, 84, width / 2, panel_top)

    pdf.setFillColor(_pdf_color(theme.text_primary))
    pdf.setFont(font_name, 16)
    pdf.drawString(left_x + 12, panel_top - 22, _pdf_text(left_title, unicode_font=unicode_font))

    pdf.setFillColor(_pdf_color(theme.accent))
    pdf.drawString(right_x + 12, panel_top - 22, _pdf_text(right_title, unicode_font=unicode_font))

    y_left = panel_top - 50
    pdf.setFillColor(_pdf_color(theme.text_primary))
    for item in left_items:
        y_left = _draw_wrapped_lines(
            pdf,
            _pdf_text(item, unicode_font=unicode_font),
            x=left_x + 12,
            y=y_left,
            max_width=panel_w - 24,
            font_name=font_name,
            font_size=int(12 * scale),
            line_height=int(16 * scale),
            max_lines=_density_limit(density, low=2, medium=3, high=4),
            prefix_first="- ",
            prefix_other="  ",
        )
        y_left -= 5

    y_right = panel_top - 50
    for item in right_items:
        y_right = _draw_wrapped_lines(
            pdf,
            _pdf_text(item, unicode_font=unicode_font),
            x=right_x + 12,
            y=y_right,
            max_width=panel_w - 24,
            font_name=font_name,
            font_size=int(12 * scale),
            line_height=int(16 * scale),
            max_lines=_density_limit(density, low=2, medium=3, high=4),
            prefix_first="- ",
            prefix_other="  ",
        )
        y_right -= 5

    if slide.image_url or slide.image_prompt:
        _draw_pdf_visual_placeholder(
            pdf,
            slide=slide,
            theme=theme,
            density=density,
            x=56,
            y=32,
            width=width - 112,
            height=44,
            font_name=font_name,
            unicode_font=unicode_font,
        )


def _render_pdf_conclusion_addons(
    pdf: canvas.Canvas,
    *,
    width: float,
    slide: SlideSpec,
    theme: DeckTheme,
    density: VisualDensity,
    font_name: str,
    unicode_font: bool,
) -> None:
    """Conclusion booster: summary cards + next steps."""
    summary_items = slide.bullets[: _density_limit(density, low=2, medium=3, high=4)]
    if not summary_items and slide.body:
        summary_items = [
            item.strip()
            for item in re.split(r"[\n;]", slide.body)
            if item.strip()
        ][:_density_limit(density, low=2, medium=3, high=4)]
    if not summary_items:
        summary_items = ["Summary 1", "Summary 2", "Summary 3"]

    card_w = (width - 56 - 56 - 24) / 3
    y = 58
    for idx, item in enumerate(summary_items):
        x = 56 + idx * (card_w + 12)
        pdf.setFillColor(_pdf_color(theme.panel_bg))
        pdf.setStrokeColor(_pdf_color(theme.panel_border))
        pdf.roundRect(x, y, card_w, 56, 8, fill=1, stroke=1)
        pdf.setFillColor(_pdf_color(theme.text_primary))
        _draw_wrapped_lines(
            pdf,
            _pdf_text(item, unicode_font=unicode_font),
            x=x + 10,
            y=y + 35,
            max_width=card_w - 18,
            font_name=font_name,
            font_size=10,
            line_height=12,
            max_lines=3,
        )

    next_steps = slide.key_message or slide.section or "Next steps"
    pdf.setFillColor(_pdf_color(theme.text_secondary))
    pdf.setFont(font_name, 10)
    pdf.drawRightString(width - 56, 42, _pdf_text(f"Next: {next_steps}", unicode_font=unicode_font))


def render_pdf_from_spec(path: Path, spec_json: dict) -> None:
    spec = PresentationSpec.model_validate(spec_json)
    theme = _select_deck_theme(spec, render_format="pdf")
    logger.debug("renderer.theme.applied format=pdf theme_variant=%s", theme.name)

    font_regular, font_bold, unicode_font = _register_pdf_fonts()
    pdf = canvas.Canvas(str(path), pagesize=A4)
    width, height = A4

    title_font = choose_pdf_font(
        spec.title,
        "title",
        regular_font=font_regular,
        bold_font=font_bold,
    )
    pdf.setTitle(_pdf_text(spec.title, unicode_font=unicode_font))
    pdf.setFont(title_font, 12)

    for idx, slide in enumerate(spec.slides):
        template = _select_template(slide, render_format="pdf")
        density = _select_visual_density(slide, template=template, render_format="pdf")
        logger.debug(
            "renderer.template.applied format=pdf slide_id=%s template=%s",
            slide.id,
            template,
        )
        logger.debug(
            "renderer.layout.applied format=pdf slide_id=%s theme_variant=%s layout=%s density=%s",
            slide.id,
            theme.name,
            template,
            density,
        )

        _paint_pdf_page_background(pdf, width=width, height=height, theme=theme)

        if template == "hero_minimal":
            _render_pdf_hero_minimal(
                pdf,
                width=width,
                height=height,
                slide=slide,
                theme=theme,
                density=density,
                font_name=font_regular,
                unicode_font=unicode_font,
            )
        elif template == "agenda_clean":
            _render_pdf_agenda_clean(
                pdf,
                width=width,
                height=height,
                slide=slide,
                theme=theme,
                density=density,
                font_name=font_regular,
                unicode_font=unicode_font,
            )
        elif template == "content_two_column":
            _render_pdf_content_two_column(
                pdf,
                width=width,
                height=height,
                slide=slide,
                theme=theme,
                density=density,
                font_name=font_regular,
                unicode_font=unicode_font,
            )
        elif template == "kpi_cards":
            _render_pdf_kpi_cards(
                pdf,
                width=width,
                height=height,
                slide=slide,
                theme=theme,
                density=density,
                font_name=font_regular,
                unicode_font=unicode_font,
            )
        elif template == "timeline_process":
            _render_pdf_timeline_process(
                pdf,
                width=width,
                height=height,
                slide=slide,
                theme=theme,
                density=density,
                font_name=font_regular,
                unicode_font=unicode_font,
            )
        elif template == "quote_focus":
            _render_pdf_quote_focus(
                pdf,
                width=width,
                height=height,
                slide=slide,
                theme=theme,
                density=density,
                font_name=font_regular,
                unicode_font=unicode_font,
            )
        elif template == "comparison_split":
            _render_pdf_comparison_split(
                pdf,
                width=width,
                height=height,
                slide=slide,
                theme=theme,
                density=density,
                font_name=font_regular,
                unicode_font=unicode_font,
            )
        else:
            _render_pdf_infographic_visual(
                pdf,
                width=width,
                height=height,
                slide=slide,
                theme=theme,
                density=density,
                font_name=font_regular,
                unicode_font=unicode_font,
            )

        if slide.type == "conclusion" and template not in {"quote_focus", "hero_minimal"}:
            _render_pdf_conclusion_addons(
                pdf,
                width=width,
                slide=slide,
                theme=theme,
                density=density,
                font_name=font_regular,
                unicode_font=unicode_font,
            )

        if idx < len(spec.slides) - 1:
            pdf.showPage()

    pdf.save()


def _pptx_set_text(
    tf,
    lines: list[str],
    *,
    base_size: int,
    bold_first: bool = False,
    rgb: tuple[int, int, int] | None = None,
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    tf.clear()
    first = tf.paragraphs[0]
    first.text = lines[0] if lines else ""
    first.font.size = Pt(base_size)
    first.font.bold = bold_first
    if rgb is not None:
        first.font.color.rgb = RGBColor(*rgb)

    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(max(12, base_size - 2))
        if rgb is not None:
            p.font.color.rgb = RGBColor(*rgb)


def _pptx_set_notes(slide: object, speaker_notes: str | None) -> None:
    if not speaker_notes:
        return
    try:
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.clear()
        notes_frame.text = speaker_notes
    except Exception:
        # TODO: if this becomes unstable for certain templates, handle notes via custom XML.
        logger.exception("renderer.notes.write_failed")


def _pptx_add_visual_placeholder(
    slide: object,
    *,
    slide_spec: SlideSpec,
    theme: DeckTheme,
    density: VisualDensity,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches

    hint = slide_spec.visual_hint or slide_spec.image_prompt or slide_spec.chart_hint
    if not hint and not slide_spec.image_url:
        return

    image_bytes = _read_image_bytes(slide_spec.image_url)
    if image_bytes:
        try:
            slide.shapes.add_picture(
                io.BytesIO(image_bytes),
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(height),
            )
            logger.debug("renderer.real_image.used format=pptx slide_id=%s", slide_spec.id)
            logger.debug(
                "renderer.visual_block.created format=pptx slide_id=%s kind=real_image",
                slide_spec.id,
            )
            return
        except Exception:
            logger.exception("renderer.image.decode_failed format=pptx slide_id=%s", slide_spec.id)

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.fill.solid()
    panel_r, panel_g, panel_b = _rgb_from_hex(theme.panel_alt_bg)
    border_r, border_g, border_b = _rgb_from_hex(theme.panel_border)
    box.fill.fore_color.rgb = RGBColor(panel_r, panel_g, panel_b)
    box.line.color.rgb = RGBColor(border_r, border_g, border_b)

    label = "Image placeholder" if slide_spec.image_prompt else "Visual block"
    lines = [label, hint or "Visual area"]
    _pptx_set_text(
        box.text_frame,
        lines,
        base_size=12 if density != "high" else 11,
        bold_first=True,
    )

    logger.debug("renderer.placeholder_image.used format=pptx slide_id=%s", slide_spec.id)
    logger.debug(
        "renderer.visual_block.created format=pptx slide_id=%s kind=%s",
        slide_spec.id,
        label,
    )


def _pptx_apply_conclusion_addons(
    slide: object,
    *,
    slide_spec: SlideSpec,
    theme: DeckTheme,
    density: VisualDensity,
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches

    item_limit = _density_limit(density, low=2, medium=3, high=4)
    items = slide_spec.bullets[:item_limit]
    if not items and slide_spec.body:
        items = [item.strip() for item in re.split(r"[\n;]", slide_spec.body) if item.strip()][:item_limit]
    if not items:
        items = ["Summary 1", "Summary 2", "Summary 3"]

    card_w = 3.7
    for idx, item in enumerate(items):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.7 + idx * (card_w + 0.25)),
            Inches(5.6),
            Inches(card_w),
            Inches(1.2),
        )
        card.fill.solid()
        panel_r, panel_g, panel_b = _rgb_from_hex(theme.panel_bg)
        border_r, border_g, border_b = _rgb_from_hex(theme.panel_border)
        card.fill.fore_color.rgb = RGBColor(panel_r, panel_g, panel_b)
        card.line.color.rgb = RGBColor(border_r, border_g, border_b)
        _pptx_set_text(card.text_frame, [item], base_size=12)

    next_steps = slide_spec.key_message or slide_spec.section or "Next steps"
    footer = slide.shapes.add_textbox(Inches(9.1), Inches(6.95), Inches(3.4), Inches(0.4))
    _pptx_set_text(footer.text_frame, [f"Next: {next_steps}"], base_size=11)


def render_pptx_from_spec(path: Path, spec_json: dict) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches

    spec = PresentationSpec.model_validate(spec_json)
    theme = _select_deck_theme(spec, render_format="pptx")
    logger.debug("renderer.theme.applied format=pptx theme_variant=%s", theme.name)

    text_primary_rgb = _rgb_from_hex(theme.text_primary)
    text_secondary_rgb = _rgb_from_hex(theme.text_secondary)
    text_inverse_rgb = _rgb_from_hex(theme.text_inverse)
    accent_rgb = _rgb_from_hex(theme.accent)
    accent_soft_rgb = _rgb_from_hex(theme.accent_soft)
    panel_bg_rgb = _rgb_from_hex(theme.panel_bg)
    panel_alt_rgb = _rgb_from_hex(theme.panel_alt_bg)
    panel_border_rgb = _rgb_from_hex(theme.panel_border)
    connector_rgb = _rgb_from_hex(theme.connector)
    page_bg_rgb = _rgb_from_hex(theme.page_bg)
    hero_bg_rgb = _rgb_from_hex(theme.hero_bg)
    rail_rgb = _rgb_from_hex(theme.rail)

    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6] if len(presentation.slide_layouts) > 6 else presentation.slide_layouts[-1]

    def add_title(slide: object, title: str, *, top: float = 0.35, size: int = 38) -> None:
        box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(12.0), Inches(1.2))
        _pptx_set_text(
            box.text_frame,
            [title],
            base_size=size,
            bold_first=True,
            rgb=text_primary_rgb,
        )

    def add_subtitle(slide: object, subtitle: str, *, top: float = 1.35) -> None:
        box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.0), Inches(0.8))
        _pptx_set_text(box.text_frame, [subtitle], base_size=18, rgb=text_secondary_rgb)

    for slide_spec in spec.slides:
        template = _select_template(slide_spec, render_format="pptx")
        density = _select_visual_density(slide_spec, template=template, render_format="pptx")
        logger.debug(
            "renderer.template.applied format=pptx slide_id=%s template=%s",
            slide_spec.id,
            template,
        )
        logger.debug(
            "renderer.layout.applied format=pptx slide_id=%s theme_variant=%s layout=%s density=%s",
            slide_spec.id,
            theme.name,
            template,
            density,
        )

        slide = presentation.slides.add_slide(blank_layout)
        _pptx_set_notes(slide, slide_spec.speaker_notes)

        # Theme-driven page background first, then template composition on top.
        bg_page = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.0),
            Inches(0.0),
            Inches(13.333),
            Inches(7.5),
        )
        bg_page.fill.solid()
        bg_page.fill.fore_color.rgb = RGBColor(*page_bg_rgb)
        bg_page.line.fill.background()

        if template == "hero_minimal":
            bg = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(0.0),
                Inches(0.0),
                Inches(13.333),
                Inches(2.1),
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(*hero_bg_rgb)
            bg.line.fill.background()

            accent = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(0.5),
                Inches(1.4),
                Inches(0.14),
                Inches(3.8),
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = RGBColor(*accent_rgb)
            accent.line.fill.background()

            add_title(
                slide,
                slide_spec.title,
                top=1.5,
                size=52 if density == "low" else (48 if density == "medium" else 44),
            )
            if slide_spec.subtitle:
                add_subtitle(slide, slide_spec.subtitle, top=3.8)
            if slide_spec.key_message:
                msg = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(8.8), Inches(0.9))
                _pptx_set_text(msg.text_frame, [slide_spec.key_message], base_size=18, rgb=text_secondary_rgb)
            _pptx_add_visual_placeholder(
                slide,
                slide_spec=slide_spec,
                theme=theme,
                density=density,
                left=9.7,
                top=4.6,
                width=2.9,
                height=1.8,
            )

        elif template == "agenda_clean":
            add_title(slide, slide_spec.title, top=0.35, size=38)
            if slide_spec.subtitle:
                add_subtitle(slide, slide_spec.subtitle, top=1.1)

            rail = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(0.95),
                Inches(2.0),
                Inches(0.04),
                Inches(4.0),
            )
            rail.fill.solid()
            rail.fill.fore_color.rgb = RGBColor(*rail_rgb)
            rail.line.fill.background()

            items = slide_spec.bullets[: _density_limit(density, low=4, medium=6, high=7)]
            if not items and slide_spec.body:
                items = [item.strip() for item in re.split(r"[\n;]", slide_spec.body) if item.strip()][
                    : _density_limit(density, low=4, medium=6, high=7)
                ]
            if not items:
                items = ["Agenda item 1", "Agenda item 2", "Agenda item 3"]

            for idx, item in enumerate(items):
                y = 2.0 + idx * 0.62
                badge = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.OVAL,
                    Inches(0.82),
                    Inches(y),
                    Inches(0.3),
                    Inches(0.3),
                )
                badge.fill.solid()
                badge.fill.fore_color.rgb = RGBColor(*accent_rgb)
                badge.line.fill.background()
                _pptx_set_text(
                    badge.text_frame,
                    [str(idx + 1)],
                    base_size=10,
                    bold_first=True,
                    rgb=text_inverse_rgb,
                )

                line_box = slide.shapes.add_textbox(Inches(1.3), Inches(y - 0.03), Inches(10.9), Inches(0.55))
                _pptx_set_text(
                    line_box.text_frame,
                    [item],
                    base_size=18 if density != "high" else 16,
                    rgb=text_primary_rgb,
                )

        elif template == "content_two_column":
            add_title(slide, slide_spec.title, top=0.35, size=36)
            if slide_spec.subtitle:
                add_subtitle(slide, slide_spec.subtitle, top=1.08)

            variant = _content_variant(slide_spec)
            left_col = (0.7, 1.9, 5.9, 4.6)
            right_col = (6.7, 1.9, 5.9, 4.6)
            text_col = left_col if variant == "left_text" else right_col
            visual_col = right_col if variant == "left_text" else left_col

            text_box = slide.shapes.add_textbox(
                Inches(text_col[0]),
                Inches(text_col[1]),
                Inches(text_col[2]),
                Inches(text_col[3]),
            )
            text_lines = slide_spec.bullets[:] or ([slide_spec.body] if slide_spec.body else ["Add bullets in spec"])
            text_lines = text_lines[: _density_limit(density, low=4, medium=6, high=8)]
            _pptx_set_text(
                text_box.text_frame,
                text_lines,
                base_size=20 if density != "high" else 18,
                rgb=text_primary_rgb,
            )

            _pptx_add_visual_placeholder(
                slide,
                slide_spec=slide_spec,
                theme=theme,
                density=density,
                left=visual_col[0],
                top=visual_col[1],
                width=visual_col[2],
                height=visual_col[3],
            )

        elif template == "kpi_cards":
            add_title(slide, slide_spec.title, top=0.35, size=36)
            if slide_spec.key_message:
                add_subtitle(slide, slide_spec.key_message, top=1.08)

            cards = _extract_cards(slide_spec, max_items=_density_limit(density, low=2, medium=4, high=4))
            card_w = 5.7
            card_h = 1.8
            for idx, (metric, label) in enumerate(cards[:4]):
                row = idx // 2
                col = idx % 2
                card = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                    Inches(0.8 + col * (card_w + 0.35)),
                    Inches(2.0 + row * (card_h + 0.35)),
                    Inches(card_w),
                    Inches(card_h),
                )
                card.fill.solid()
                if idx % 2 == 0:
                    card.fill.fore_color.rgb = RGBColor(*panel_bg_rgb)
                    card.line.color.rgb = RGBColor(*panel_border_rgb)
                else:
                    card.fill.fore_color.rgb = RGBColor(*panel_alt_rgb)
                    card.line.color.rgb = RGBColor(*panel_border_rgb)
                _pptx_set_text(
                    card.text_frame,
                    [metric, label],
                    base_size=30 if density != "low" else 34,
                    bold_first=True,
                    rgb=text_primary_rgb,
                )

        elif template == "timeline_process":
            add_title(slide, slide_spec.title, top=0.35, size=36)
            if slide_spec.subtitle:
                add_subtitle(slide, slide_spec.subtitle, top=1.08)

            steps = _extract_timeline_steps(
                slide_spec,
                max_items=_density_limit(density, low=4, medium=6, high=8),
            )
            count = len(steps)
            step_w = min(2.0, (11.2 - max(0, count - 1) * 0.24) / max(1, count))

            for idx, step in enumerate(steps):
                left = 1.0 + idx * (step_w + 0.24)
                step_box = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                    Inches(left),
                    Inches(3.0),
                    Inches(step_w),
                    Inches(1.45),
                )
                step_box.fill.solid()
                step_box.fill.fore_color.rgb = RGBColor(*panel_bg_rgb)
                step_box.line.color.rgb = RGBColor(*panel_border_rgb)
                _pptx_set_text(
                    step_box.text_frame,
                    [f"{idx + 1}", step],
                    base_size=16 if density != "high" else 14,
                    bold_first=True,
                    rgb=text_primary_rgb,
                )

                if idx < count - 1:
                    connector = slide.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                        Inches(left + step_w),
                        Inches(3.65),
                        Inches(0.24),
                        Inches(0.06),
                    )
                    connector.fill.solid()
                    connector.fill.fore_color.rgb = RGBColor(*connector_rgb)
                    connector.line.fill.background()

            _pptx_add_visual_placeholder(
                slide,
                slide_spec=slide_spec,
                theme=theme,
                density=density,
                left=9.7,
                top=5.0,
                width=2.8,
                height=1.6,
            )

        elif template == "quote_focus":
            logger.debug("renderer.quote_layout.used format=pptx slide_id=%s", slide_spec.id)

            bg = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(0.0),
                Inches(0.0),
                Inches(13.333),
                Inches(7.5),
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(*page_bg_rgb)
            bg.line.fill.background()

            panel = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(0.6),
                Inches(0.6),
                Inches(12.1),
                Inches(6.2),
            )
            panel.fill.solid()
            panel.fill.fore_color.rgb = RGBColor(*panel_alt_rgb)
            panel.line.fill.background()

            quote = _extract_quote_text(slide_spec)
            quote_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.4), Inches(10.8), Inches(3.8))
            _pptx_set_text(
                quote_box.text_frame,
                [f"\"{quote}\""],
                base_size=42 if density == "low" else (40 if density == "medium" else 34),
                rgb=text_primary_rgb,
            )

            if slide_spec.subtitle:
                author = slide.shapes.add_textbox(Inches(1.2), Inches(5.4), Inches(8.0), Inches(0.6))
                _pptx_set_text(author.text_frame, [slide_spec.subtitle], base_size=16, rgb=text_secondary_rgb)

            _pptx_add_visual_placeholder(
                slide,
                slide_spec=slide_spec,
                theme=theme,
                density=density,
                left=9.8,
                top=5.2,
                width=2.6,
                height=1.2,
            )

        elif template == "comparison_split":
            logger.debug("renderer.comparison_layout.used format=pptx slide_id=%s", slide_spec.id)

            add_title(slide, slide_spec.title, top=0.35, size=36)
            left_title, left_items, right_title, right_items = _extract_comparison_columns(slide_spec)

            left_panel = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(0.7),
                Inches(1.6),
                Inches(5.9),
                Inches(4.9),
            )
            left_panel.fill.solid()
            left_panel.fill.fore_color.rgb = RGBColor(*panel_bg_rgb)
            left_panel.line.color.rgb = RGBColor(*panel_border_rgb)
            _pptx_set_text(
                left_panel.text_frame,
                [left_title, *left_items[: _density_limit(density, low=3, medium=4, high=5)]],
                base_size=20 if density != "high" else 18,
                bold_first=True,
                rgb=text_primary_rgb,
            )

            right_panel = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(6.75),
                Inches(1.6),
                Inches(5.9),
                Inches(4.9),
            )
            right_panel.fill.solid()
            right_panel.fill.fore_color.rgb = RGBColor(*panel_alt_rgb)
            right_panel.line.color.rgb = RGBColor(*panel_border_rgb)
            _pptx_set_text(
                right_panel.text_frame,
                [right_title, *right_items[: _density_limit(density, low=3, medium=4, high=5)]],
                base_size=20 if density != "high" else 18,
                bold_first=True,
                rgb=text_primary_rgb,
            )

            divider = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(6.63),
                Inches(1.7),
                Inches(0.05),
                Inches(4.7),
            )
            divider.fill.solid()
            divider.fill.fore_color.rgb = RGBColor(*connector_rgb)
            divider.line.fill.background()

            _pptx_add_visual_placeholder(
                slide,
                slide_spec=slide_spec,
                theme=theme,
                density=density,
                left=0.7,
                top=6.6,
                width=12.0,
                height=0.7,
            )

        else:  # infographic_visual
            add_title(slide, slide_spec.title, top=0.35, size=36)
            center = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(4.5),
                Inches(2.4),
                Inches(4.0),
                Inches(2.0),
            )
            center.fill.solid()
            center.fill.fore_color.rgb = RGBColor(*accent_soft_rgb)
            center.line.color.rgb = RGBColor(*panel_border_rgb)
            center_text = slide_spec.key_message or slide_spec.body or "Core concept"
            _pptx_set_text(center.text_frame, [center_text], base_size=16, rgb=text_primary_rgb)

            nodes = _extract_infographic_nodes(
                slide_spec,
                max_items=_density_limit(density, low=4, medium=6, high=6),
            )
            positions = [
                (1.0, 1.8),
                (9.4, 1.8),
                (1.0, 4.8),
                (9.4, 4.8),
                (4.8, 1.0),
                (4.8, 5.8),
            ]
            for idx, text in enumerate(nodes):
                left, top = positions[idx]
                node = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                    Inches(left),
                    Inches(top),
                    Inches(2.7),
                    Inches(1.0),
                )
                node.fill.solid()
                node.fill.fore_color.rgb = RGBColor(*panel_bg_rgb)
                node.line.color.rgb = RGBColor(*panel_border_rgb)
                _pptx_set_text(node.text_frame, [text], base_size=12, rgb=text_primary_rgb)

            _pptx_add_visual_placeholder(
                slide,
                slide_spec=slide_spec,
                theme=theme,
                density=density,
                left=0.9,
                top=6.2,
                width=3.2,
                height=1.0,
            )

        if slide_spec.type == "conclusion":
            _pptx_apply_conclusion_addons(
                slide,
                slide_spec=slide_spec,
                theme=theme,
                density=density,
            )

    presentation.save(str(path))


def render_placeholder_pdf(path: Path, job_id: str, prompt: str) -> None:
    pdf = canvas.Canvas(str(path), pagesize=A4)
    font_regular, font_bold, unicode_font = _register_pdf_fonts()
    _, height = A4
    title_text = _pdf_text(f"Presentation {job_id}", unicode_font=unicode_font)
    body_text = _pdf_text("Placeholder content (spec unavailable).", unicode_font=unicode_font)
    prompt_text = _pdf_text(prompt[:300], unicode_font=unicode_font)

    pdf.setFont(
        choose_pdf_font(title_text, "title", regular_font=font_regular, bold_font=font_bold),
        18,
    )
    pdf.drawString(56, height - 90, title_text)
    pdf.setFont(
        choose_pdf_font(body_text, "body", regular_font=font_regular, bold_font=font_bold),
        12,
    )
    pdf.drawString(56, height - 130, body_text)
    pdf.drawString(56, height - 155, prompt_text)
    pdf.save()


def render_placeholder_pptx(path: Path, job_id: str, prompt: str) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = f"Presentation {job_id}"
    for placeholder in slide.placeholders:
        if placeholder != slide.shapes.title and hasattr(placeholder, "text"):
            placeholder.text = f"Placeholder content\\n{prompt[:240]}"
            break
    presentation.save(str(path))
